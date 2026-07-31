"""Generates the 17-slide defense presentation (.pptx) for the
Intelligent Heavy Mobile Equipment Predictive Maintenance System.

Run with:
    python scripts/generate_presentation.py

Produces docs/soutenance_HME.pptx — a self-contained, modern 16:9 deck
(Microsoft / Tesla / IBM Watson / Siemens / Power BI inspired: dark navy,
white, gray, orange), reusing the real chart images already generated in
docs/report_assets/ (run scripts/generate_report_assets.py first if that
folder is empty or stale).
"""

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

BASE_DIR = Path(__file__).resolve().parent.parent
ASSETS_DIR = BASE_DIR / "docs" / "report_assets"
OUTPUT_PATH = BASE_DIR / "docs" / "soutenance_HME.pptx"

# ---------------------------------------------------------------------------
# Palette (Microsoft / Tesla / IBM Watson / Siemens / Power BI inspired)
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x0B, 0x1F, 0x3A)
NAVY_2 = RGBColor(0x12, 0x2A, 0x4D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF4, 0xF6, 0xFA)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x64, 0x74, 0x8B)
GRAY_DARK = RGBColor(0x33, 0x41, 0x55)
ORANGE = RGBColor(0xF9, 0x73, 0x16)
ORANGE_DARK = RGBColor(0xC2, 0x54, 0x0A)
BLUE_ACCENT = RGBColor(0x25, 0x63, 0xEB)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)

FONT_HEAD = "Segoe UI Semibold"
FONT_BODY = "Segoe UI"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------

def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def no_line(shape):
    shape.line.fill.background()


def add_rect(slide, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE, line_color=None, shadow=False):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line_color is None:
        no_line(sp)
    else:
        sp.line.color.rgb = line_color
        sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    return sp


def add_text(slide, x, y, w, h, text, size=18, color=GRAY_DARK, bold=False, italic=False,
             font=FONT_BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0,
             wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font
        run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, size=15, color=GRAY_DARK, marker_color=ORANGE,
                 gap=Pt(10), font=FONT_BODY, bold_lead=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = gap
        p.line_spacing = 1.05
        marker = p.add_run()
        marker.text = "▸  "
        marker.font.size = Pt(size)
        marker.font.bold = True
        marker.font.color.rgb = marker_color
        marker.font.name = font
        if isinstance(item, tuple):
            lead, rest = item
            r1 = p.add_run()
            r1.text = lead
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = color
            r1.font.name = font
            r2 = p.add_run()
            r2.text = rest
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
            r2.font.name = font
        else:
            r = p.add_run()
            r.text = item
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.name = font
    return tb


def add_icon_badge(slide, x, y, size, symbol, bg=NAVY, fg=WHITE, shape=MSO_SHAPE.OVAL):
    sp = add_rect(slide, x, y, size, size, bg, shape=shape)
    tf = sp.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = symbol
    run.font.size = Pt(int(size.inches * 26))
    run.font.color.rgb = fg
    return sp


def header_band(slide, kicker, title, dark=False):
    """Top header used on every content slide: small orange kicker + big title."""

    text_color = WHITE if dark else NAVY
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.09), ORANGE)
    add_text(slide, Inches(0.7), Inches(0.42), Inches(9), Inches(0.35), kicker.upper(),
              size=13, color=ORANGE, bold=True, font=FONT_BODY)
    add_text(slide, Inches(0.7), Inches(0.72), Inches(11.9), Inches(0.7), title,
              size=30, color=text_color, bold=True, font=FONT_HEAD)
    add_rect(slide, Inches(0.72), Inches(1.42), Inches(0.6), Pt(3), ORANGE)


def footer(slide, number, dark=False):
    color = RGBColor(0x9C, 0xA8, 0xC0) if dark else GRAY
    add_text(slide, Inches(0.7), Inches(7.14), Inches(7), Inches(0.3),
              "Intelligent HME Predictive Maintenance System", size=9.5, color=color, italic=True)
    add_text(slide, Inches(12.4), Inches(7.14), Inches(0.6), Inches(0.3), str(number),
              size=9.5, color=color, align=PP_ALIGN.RIGHT, bold=True)


def picture_framed(slide, path, x, y, w, border=BORDER):
    pic = slide.shapes.add_picture(str(path), x, y, width=w)
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, pic.width, pic.height)
    rect.fill.background()
    rect.line.color.rgb = border
    rect.line.width = Pt(1)
    rect.shadow.inherit = False
    return pic


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def slide_01_title(prs):
    s = new_slide(prs)
    set_bg(s, NAVY)
    add_rect(s, Inches(0), Inches(6.55), SLIDE_W, Inches(0.06), ORANGE)
    add_rect(s, Inches(0.9), Inches(1.55), Inches(0.55), Pt(4), ORANGE)
    add_text(s, Inches(0.9), Inches(1.85), Inches(11.5), Inches(0.5),
              "MÉMOIRE DE MASTER — INFORMATIQUE / GÉNIE LOGICIEL / SCIENCE DES DONNÉES",
              size=15, color=ORANGE, bold=True)
    add_text(s, Inches(0.85), Inches(2.35), Inches(11.6), Inches(1.9),
              "Intelligent Heavy Mobile\nEquipment Predictive\nMaintenance System",
              size=42, color=WHITE, bold=True, font=FONT_HEAD, line_spacing=1.03)
    add_text(s, Inches(0.9), Inches(4.55), Inches(10.8), Inches(0.9),
              "Un système intelligent de maintenance prédictive pour équipements miniers lourds, "
              "fondé sur le Machine Learning, l'explicabilité (XAI) et les pratiques MLOps.",
              size=16, color=RGBColor(0xC7, 0xD2, 0xE3), italic=True, line_spacing=1.2)
    add_text(s, Inches(0.9), Inches(6.65), Inches(6), Inches(0.4),
              "[Nom et Prénom de l'étudiant·e]  •  Encadreur : [ ]", size=13, color=WHITE, bold=True)
    add_text(s, Inches(10.2), Inches(6.65), Inches(2.3), Inches(0.4),
              "2025 — 2026", size=13, color=RGBColor(0xC7, 0xD2, 0xE3), align=PP_ALIGN.RIGHT)


def slide_02_contexte(prs):
    s = new_slide(prs)
    set_bg(s, LIGHT_BG)
    header_band(s, "Contexte", "Un enjeu industriel majeur pour les sites miniers")
    icons = [
        ("🚛", "Flottes lourdes", "Excavatrices, chargeuses, foreuses, camions, niveleuses, bulldozers"),
        ("⛔", "Arrêts non planifiés", "Un équipement immobilisé rompt toute la chaîne de production"),
        ("💰", "Coûts considérables", "Pièces en urgence, mobilisation d'équipes, perte de production"),
    ]
    card_w = Inches(3.75)
    gap = Inches(0.35)
    start_x = Inches(0.7)
    for i, (icon, title, desc) in enumerate(icons):
        x = start_x + i * (card_w + gap)
        card = add_rect(s, x, Inches(2.1), card_w, Inches(3.6), CARD_BG, line_color=BORDER)
        add_icon_badge(s, x + Inches(1.375), Inches(2.45), Inches(1.0), icon, bg=NAVY, fg=WHITE)
        add_text(s, x + Inches(0.25), Inches(3.65), card_w - Inches(0.5), Inches(0.5), title,
                  size=17, bold=True, color=NAVY, align=PP_ALIGN.CENTER, font=FONT_HEAD)
        add_text(s, x + Inches(0.3), Inches(4.15), card_w - Inches(0.6), Inches(1.4), desc,
                  size=12.5, color=GRAY_DARK, align=PP_ALIGN.CENTER, line_spacing=1.15)
    add_text(s, Inches(0.7), Inches(6.0), Inches(11.9), Inches(0.9),
              "Sites concernés : Kamoa Copper, Tenke Fungurume Mining, MMG Kinsevere, Ivanhoe Mines, "
              "Glencore, CMOC…", size=13, color=GRAY, italic=True, align=PP_ALIGN.CENTER)
    footer(s, 2)


def slide_03_problematique(prs):
    s = new_slide(prs)
    set_bg(s, LIGHT_BG)
    header_band(s, "Problématique", "Deux stratégies de maintenance, deux limites")

    left = add_rect(s, Inches(0.7), Inches(2.05), Inches(5.55), Inches(3.4), CARD_BG, line_color=BORDER)
    add_rect(s, Inches(0.7), Inches(2.05), Inches(5.55), Inches(0.55), NAVY)
    add_text(s, Inches(0.7), Inches(2.05), Inches(5.55), Inches(0.55), "Maintenance corrective",
              size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
    add_bullets(s, Inches(1.05), Inches(2.85), Inches(4.9), Inches(2.4), [
        "Intervention après la panne",
        "Simple à mettre en œuvre",
        "Arrêts non planifiés coûteux",
        "Risques pour la sécurité",
    ], size=14)

    right = add_rect(s, Inches(7.05), Inches(2.05), Inches(5.55), Inches(3.4), CARD_BG, line_color=BORDER)
    add_rect(s, Inches(7.05), Inches(2.05), Inches(5.55), Inches(0.55), GRAY_DARK)
    add_text(s, Inches(7.05), Inches(2.05), Inches(5.55), Inches(0.55), "Maintenance préventive systématique",
              size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
    add_bullets(s, Inches(7.4), Inches(2.85), Inches(4.9), Inches(2.4), [
        "Intervalles fixes (ex. 500h)",
        "Réduit le risque de panne brutale",
        "Remplace des pièces encore bonnes",
        "Surcoût important",
    ], size=14)

    add_icon_badge(s, Inches(6.13), Inches(3.35), Inches(1.05), "VS", bg=ORANGE, fg=WHITE,
                    shape=MSO_SHAPE.OVAL)
    box = add_rect(s, Inches(0.7), Inches(5.75), Inches(11.9), Inches(1.0), NAVY,
                    shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, Inches(1.1), Inches(5.75), Inches(11.1), Inches(1.0),
              "Ni l'une ni l'autre ne tient compte de l'état réel de chaque équipement, mesuré en continu.",
              size=15, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 3)


def slide_04_objectifs(prs):
    s = new_slide(prs)
    set_bg(s, LIGHT_BG)
    header_band(s, "Objectifs", "Ce que le système doit accomplir")
    items = [
        ("🎯", "Prédire", "le risque de panne à 7 jours, par équipement et pour toute la flotte"),
        ("📉", "Réduire", "les temps d'arrêt non planifiés en identifiant les équipements critiques"),
        ("🧠", "Expliquer", "chaque prédiction avec l'IA explicable (SHAP)"),
        ("📊", "Visualiser", "un tableau de bord interactif (disponibilité, MTBF, MTTR)"),
        ("🔄", "Automatiser", "tout le cycle de vie du modèle (entraînement, dérive, réentraînement)"),
    ]
    card_w = Inches(2.24)
    gap = Inches(0.18)
    start_x = Inches(0.7)
    for i, (icon, title, desc) in enumerate(items):
        x = start_x + i * (card_w + gap)
        add_rect(s, x, Inches(2.15), card_w, Inches(4.4), CARD_BG, line_color=BORDER,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_icon_badge(s, x + card_w / 2 - Inches(0.45), Inches(2.45), Inches(0.9), icon, bg=ORANGE, fg=WHITE)
        add_text(s, x + Inches(0.12), Inches(3.55), card_w - Inches(0.24), Inches(0.4), title,
                  size=15.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER, font=FONT_HEAD)
        add_text(s, x + Inches(0.15), Inches(4.0), card_w - Inches(0.3), Inches(2.4), desc,
                  size=11, color=GRAY_DARK, align=PP_ALIGN.CENTER, line_spacing=1.15)
    footer(s, 4)


def slide_05_dataset(prs):
    s = new_slide(prs)
    set_bg(s, LIGHT_BG)
    header_band(s, "Données", "Le jeu de données HME_Downtime")
    stats = [("10 000", "lectures capteurs"), ("300", "équipements"), ("6", "types d'équipement"),
             ("9", "colonnes"), ("24,6 %", "taux de panne (cible)")]
    card_w = Inches(2.2)
    gap = Inches(0.22)
    start_x = Inches(0.7)
    for i, (val, label) in enumerate(stats):
        x = start_x + i * (card_w + gap)
        add_rect(s, x, Inches(2.2), card_w, Inches(1.7), NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text(s, x, Inches(2.4), card_w, Inches(0.75), val, size=28, bold=True, color=ORANGE,
                  align=PP_ALIGN.CENTER, font=FONT_HEAD)
        add_text(s, x + Inches(0.1), Inches(3.15), card_w - Inches(0.2), Inches(0.6), label,
                  size=11.5, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(4.3), Inches(11.9), Inches(0.4), "Schéma des variables",
              size=15, bold=True, color=NAVY, font=FONT_HEAD)
    add_bullets(s, Inches(0.7), Inches(4.8), Inches(5.8), Inches(2.2), [
        ("Identifiants — ", "EquipmentID, EquipmentType, Timestamp"),
        ("Capteurs — ", "OperatingHours, EngineTemp, HydraulicPressure, Vibration"),
        ("Contexte — ", "FailureMode (dernier mode de panne observé)"),
    ], size=13.5)
    add_bullets(s, Inches(6.7), Inches(4.8), Inches(5.9), Inches(2.2), [
        ("Cible — ", "FailureWithin7Days (0 = pas de panne, 1 = panne sous 7 jours)"),
        ("Qualité — ", "0 valeur manquante, 1 doublon supprimé"),
    ], size=13.5)
    footer(s, 5)


def slide_06_feature_engineering(prs):
    s = new_slide(prs)
    set_bg(s, LIGHT_BG)
    header_band(s, "Préparation des données", "Ingénierie de caractéristiques (Feature Engineering)")
    add_icon_badge(s, Inches(10.6), Inches(2.15), Inches(2.0), "57", bg=ORANGE, fg=WHITE,
                    shape=MSO_SHAPE.OVAL)
    add_text(s, Inches(10.1), Inches(4.25), Inches(3.0), Inches(0.5), "variables dérivées",
              size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    families = [
        ("⏱", "Temporelles", "Moyennes/écarts-types glissants, tendances, valeurs décalées"),
        ("⚠", "Scores de risque", "Risque thermique, pression, vibration, santé composite"),
        ("📆", "Usage & cycle de vie", "Âge relatif, durée de vie résiduelle estimée"),
        ("🔗", "Interactions", "Produits/ratios entre capteurs, z-scores par type"),
        ("🏷", "Encodages", "Type d'équipement, fréquence du mode de panne, cyclique temporel"),
    ]
    y = Inches(2.15)
    for icon, title, desc in families:
        add_icon_badge(s, Inches(0.7), y, Inches(0.65), icon, bg=NAVY, fg=WHITE)
        add_text(s, Inches(1.55), y - Inches(0.02), Inches(3.0), Inches(0.4), title, size=14.5,
                  bold=True, color=NAVY, font=FONT_HEAD)
        add_text(s, Inches(1.55), y + Inches(0.36), Inches(7.9), Inches(0.6), desc, size=12,
                  color=GRAY_DARK)
        y += Inches(0.92)
    footer(s, 6)


def slide_07_automl(prs):
    s = new_slide(prs)
    set_bg(s, LIGHT_BG)
    header_band(s, "Modélisation", "Pipeline AutoML + optimisation Optuna")
    models = ["Random\nForest", "Extra\nTrees", "Gradient\nBoosting", "HistGB", "XGBoost", "LightGBM", "CatBoost"]
    card_w = Inches(1.62)
    gap = Inches(0.12)
    start_x = Inches(0.7)
    for i, m in enumerate(models):
        x = start_x + i * (card_w + gap)
        add_rect(s, x, Inches(2.1), card_w, Inches(1.15), CARD_BG, line_color=BLUE_ACCENT,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text(s, x, Inches(2.1), card_w, Inches(1.15), m, size=11.5, bold=True, color=NAVY,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.7), Inches(3.45), Inches(11.9), Inches(0.4),
              "▼  Comparaison automatisée (score composite F1 / ROC AUC)", size=13, color=GRAY,
              align=PP_ALIGN.CENTER, italic=True)
    box = add_rect(s, Inches(2.9), Inches(4.0), Inches(7.5), Inches(1.0), ORANGE,
                    shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, Inches(2.9), Inches(4.0), Inches(7.5), Inches(1.0),
              "Top 3 modèles  →  Optuna : 30 essais d'optimisation bayésienne par modèle",
              size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.7), Inches(5.3), Inches(11.9), Inches(0.4),
              "▼  Sélection automatique du meilleur modèle", size=13, color=GRAY,
              align=PP_ALIGN.CENTER, italic=True)
    add_rect(s, Inches(4.4), Inches(5.85), Inches(4.5), Inches(0.95), NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, Inches(4.4), Inches(5.85), Inches(4.5), Inches(0.95), "Modèle actif : CatBoost (optimisé)",
              size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 7)


def slide_08_architecture(prs):
    s = new_slide(prs)
    set_bg(s, LIGHT_BG)
    header_band(s, "Architecture", "Une architecture logicielle en couches")
    layers = [
        ("Interface Web", "Flask · Bootstrap 5 · Chart.js / ApexCharts", NAVY),
        ("API REST", "FastAPI · JWT · Swagger", BLUE_ACCENT),
        ("Services métier", "Prediction · Equipment · Dashboard · Explainability", GRAY_DARK),
        ("Machine Learning", "AutoML · Optuna · SHAP · Registre de modèles", ORANGE_DARK),
        ("Données", "SQLAlchemy · SQLite (12 tables)", RGBColor(0x33, 0x41, 0x55)),
    ]
    y = Inches(2.05)
    h = Inches(0.92)
    for name, desc, color in layers:
        add_rect(s, Inches(1.3), y, Inches(10.7), h, color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text(s, Inches(1.7), y, Inches(3.6), h, name, size=16, bold=True, color=WHITE,
                  anchor=MSO_ANCHOR.MIDDLE, font=FONT_HEAD)
        add_text(s, Inches(5.3), y, Inches(6.5), h, desc, size=12.5,
                  color=RGBColor(0xE2, 0xE8, 0xF0), anchor=MSO_ANCHOR.MIDDLE)
        y += h + Inches(0.13)
    footer(s, 8)


def slide_09_resultats(prs):
    s = new_slide(prs)
    set_bg(s, LIGHT_BG)
    header_band(s, "Résultats", "Classement des modèles optimisés")
    rows = [
        ("Modèle", "Accuracy", "F1", "ROC AUC", "MCC"),
        ("CatBoost (optimisé) ★", "0,739", "0,078", "0,511", "0,022"),
        ("XGBoost (optimisé)", "0,746", "0,015", "0,518", "-0,022"),
        ("Extra Trees (optimisé)", "0,747", "0,027", "0,485", "0,001"),
    ]
    x, y, w, h = Inches(1.4), Inches(2.2), Inches(10.5), Inches(2.2)
    table_shape = s.shapes.add_table(len(rows), len(rows[0]), x, y, w, h)
    table = table_shape.table
    widths = [Inches(3.5), Inches(2.1), Inches(1.8), Inches(1.6), Inches(1.5)]
    for i, wcol in enumerate(widths):
        table.columns[i].width = wcol
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            run = p.runs[0]
            run.font.size = Pt(13)
            run.font.name = FONT_BODY
            run.font.bold = (r == 0)
            run.font.color.rgb = WHITE if r == 0 else (ORANGE_DARK if r == 1 else GRAY_DARK)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if r == 0 else (RGBColor(0xFF, 0xF1, 0xE6) if r == 1 else CARD_BG)
    add_rect(s, Inches(1.4), Inches(4.75), Inches(10.5), Inches(1.5), NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, Inches(1.8), Inches(4.9), Inches(9.7), Inches(1.2),
              "Modèle retenu : CatBoost optimisé — score composite le plus élevé (0,2945)",
              size=16, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 9)


def slide_10_decouverte(prs):
    s = new_slide(prs)
    set_bg(s, NAVY)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.09), ORANGE)
    add_text(s, Inches(0.7), Inches(0.42), Inches(10), Inches(0.35), "RÉSULTAT SCIENTIFIQUE CLÉ",
              size=13, color=ORANGE, bold=True)
    add_text(s, Inches(0.7), Inches(0.72), Inches(11.9), Inches(0.7),
              "Une découverte, pas un échec", size=30, color=WHITE, bold=True, font=FONT_HEAD)

    add_text(s, Inches(0.7), Inches(1.75), Inches(5.3), Inches(1.3), "ROC AUC ≈ 0,51",
              size=44, bold=True, color=ORANGE, font=FONT_HEAD)
    add_text(s, Inches(0.75), Inches(2.85), Inches(5.2), Inches(0.5), "niveau du hasard",
              size=15, italic=True, color=RGBColor(0xC7, 0xD2, 0xE3))
    add_bullets(s, Inches(0.7), Inches(3.55), Inches(5.4), Inches(3.2), [
        "Corrélations capteurs ↔ cible < 0,02 (Pearson)",
        "Information mutuelle ≈ 0",
        "Confirmé par l'analyse SHAP post-entraînement",
        "→ Limite du dataset synthétique, pas du pipeline",
    ], size=15, color=WHITE, marker_color=ORANGE)

    picture_framed(s, ASSETS_DIR / "correlation_heatmap.png", Inches(6.7), Inches(1.95), Inches(5.4))
    footer(s, 10, dark=True)


def slide_11_shap(prs):
    s = new_slide(prs)
    set_bg(s, LIGHT_BG)
    header_band(s, "Explicabilité", "Comprendre chaque prédiction (SHAP)")
    picture_framed(s, ASSETS_DIR / "shap_summary.png", Inches(0.7), Inches(2.05), Inches(5.3))
    add_bullets(s, Inches(7.65), Inches(2.2), Inches(5.1), Inches(4.5), [
        ("SHAP — ", "attribue à chaque variable une contribution signée à la prédiction"),
        ("Pourquoi TreeExplainer ? ", "explicateur exact pour les modèles à base d'arbres"),
        ("En production — ", "chaque prédiction est accompagnée d'une explication en langage naturel"),
        ("Recommandation — ", "niveau de confiance + action de maintenance suggérée"),
    ], size=13.5)
    footer(s, 11)


def slide_12_dashboard(prs):
    s = new_slide(prs)
    set_bg(s, LIGHT_BG)
    header_band(s, "Démonstration", "Tableau de bord & équipements à risque")
    kpis = [
        ("300", "Équipements", NAVY), ("2 460", "Pannes détectées", ORANGE_DARK),
        ("97,7 %", "Disponibilité", RGBColor(0x16, 0xA3, 0x4A)), ("1 068 h", "MTBF", BLUE_ACCENT),
        ("25 h", "MTTR", ORANGE), ("81", "Équip. critiques", RGBColor(0x99, 0x1B, 0x1B)),
    ]
    card_w = Inches(1.87)
    gap = Inches(0.15)
    start_x = Inches(0.7)
    for i, (val, label, color) in enumerate(kpis):
        x = start_x + i * (card_w + gap)
        add_rect(s, x, Inches(2.1), card_w, Inches(1.4), CARD_BG, line_color=BORDER,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_rect(s, x, Inches(2.1), Inches(0.08), Inches(1.4), color)
        add_text(s, x + Inches(0.15), Inches(2.25), card_w - Inches(0.3), Inches(0.6), val,
                  size=20, bold=True, color=NAVY, font=FONT_HEAD)
        add_text(s, x + Inches(0.15), Inches(2.85), card_w - Inches(0.3), Inches(0.55), label,
                  size=10.5, color=GRAY_DARK)
    add_text(s, Inches(0.7), Inches(3.85), Inches(11.9), Inches(0.4),
              "Page « Équipements à risque » — classement en temps réel", size=15, bold=True,
              color=NAVY, font=FONT_HEAD)
    add_bullets(s, Inches(0.7), Inches(4.35), Inches(11.5), Inches(2.2), [
        "Calcul de la probabilité de panne pour toute la flotte, sur demande",
        "Classement décroissant par probabilité, avec badge de risque (Faible / Moyen / Élevé / Critique)",
        "Génération automatique d'alertes pour les équipements critiques",
    ], size=13.5)
    footer(s, 12)


def slide_13_mlops(prs):
    s = new_slide(prs)
    set_bg(s, LIGHT_BG)
    header_band(s, "MLOps", "Cycle de vie automatisé du modèle")
    steps = [
        ("💾", "Registre de modèles", "Versions tracées, une seule active à la fois"),
        ("📉", "Détection de dérive", "PSI + test de Kolmogorov-Smirnov"),
        ("🔄", "Réentraînement auto", "Déclenché automatiquement si dérive détectée"),
    ]
    card_w = Inches(3.75)
    gap = Inches(0.35)
    start_x = Inches(0.7)
    for i, (icon, title, desc) in enumerate(steps):
        x = start_x + i * (card_w + gap)
        add_rect(s, x, Inches(2.2), card_w, Inches(3.2), CARD_BG, line_color=BORDER,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_icon_badge(s, x + card_w / 2 - Inches(0.5), Inches(2.55), Inches(1.0), icon, bg=NAVY, fg=WHITE)
        add_text(s, x + Inches(0.2), Inches(3.75), card_w - Inches(0.4), Inches(0.5), title, size=15,
                  bold=True, color=NAVY, align=PP_ALIGN.CENTER, font=FONT_HEAD)
        add_text(s, x + Inches(0.3), Inches(4.25), card_w - Inches(0.6), Inches(1.0), desc, size=12,
                  color=GRAY_DARK, align=PP_ALIGN.CENTER, line_spacing=1.15)
        if i < 2:
            add_text(s, x + card_w + Inches(0.02), Inches(3.55), Inches(0.32), Inches(0.6), "→",
                      size=26, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(5.75), Inches(11.9), Inches(0.6),
              "Alternative légère mais réellement fonctionnelle à un orchestrateur complet (Airflow)",
              size=13, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
    footer(s, 13)


def slide_14_limites(prs):
    s = new_slide(prs)
    set_bg(s, LIGHT_BG)
    header_band(s, "Discussion", "Limites identifiées")
    add_bullets(s, Inches(0.9), Inches(2.3), Inches(11.3), Inches(3.8), [
        ("Signal du dataset — ", "aucune corrélation exploitable entre capteurs bruts et panne "
         "dans ce jeu de données synthétique"),
        ("PredictedRemainingLife — ", "légère fuite d'information au niveau population, à corriger "
         "en production stricte"),
        ("Docker / MinIO / Airflow — ", "préparés mais non exécutés (aucun démon Docker disponible "
         "dans l'environnement de développement)"),
        ("Concept drift — ", "non implémenté ; nécessite un flux continu d'étiquettes réelles"),
    ], size=16, gap=Pt(16))
    footer(s, 14)


def slide_15_perspectives(prs):
    s = new_slide(prs)
    set_bg(s, LIGHT_BG)
    header_band(s, "Discussion", "Perspectives")
    add_bullets(s, Inches(0.9), Inches(2.3), Inches(11.3), Inches(3.8), [
        ("Données réelles — ", "reconnecter le pipeline à un flux de télémétrie réel d'un site minier"),
        ("Capteurs enrichis — ", "niveau d'huile, courant électrique, données GPS de charge utile"),
        ("Concept drift — ", "surveillance de la relation capteurs → panne dans le temps"),
        ("Infrastructure — ", "PostgreSQL, déploiement conteneurisé validé, orchestration renforcée"),
    ], size=16, gap=Pt(16))
    footer(s, 15)


def slide_16_conclusion(prs):
    s = new_slide(prs)
    set_bg(s, LIGHT_BG)
    header_band(s, "Conclusion", "Une base solide, prête pour des données réelles")
    add_bullets(s, Inches(0.9), Inches(2.2), Inches(11.3), Inches(3.6), [
        "Architecture logicielle complète et réutilisable (MVC, API REST, MLOps léger)",
        "Chaîne de traitement rigoureuse : données → features → AutoML → explicabilité",
        "La rigueur scientifique a permis de détecter une limite du dataset plutôt que de la masquer",
        "Prêt pour un déploiement réel dès qu'un signal prédictif authentique est disponible",
    ], size=16, gap=Pt(14))
    box = add_rect(s, Inches(0.9), Inches(5.75), Inches(11.3), Inches(1.0), NAVY,
                    shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, Inches(1.3), Inches(5.75), Inches(10.5), Inches(1.0),
              "Un système, une méthode, une leçon scientifique.", size=17, bold=True, color=WHITE,
              anchor=MSO_ANCHOR.MIDDLE, italic=True)
    footer(s, 16)


def slide_17_merci(prs):
    s = new_slide(prs)
    set_bg(s, NAVY)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.09), ORANGE)
    add_text(s, Inches(0), Inches(2.9), SLIDE_W, Inches(1.2), "Merci de votre attention",
              size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=FONT_HEAD)
    add_text(s, Inches(0), Inches(4.1), SLIDE_W, Inches(0.6), "Questions ?",
              size=22, color=ORANGE, align=PP_ALIGN.CENTER, italic=True)
    footer(s, 17, dark=True)


def main() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for builder in [
        slide_01_title, slide_02_contexte, slide_03_problematique, slide_04_objectifs,
        slide_05_dataset, slide_06_feature_engineering, slide_07_automl, slide_08_architecture,
        slide_09_resultats, slide_10_decouverte, slide_11_shap, slide_12_dashboard,
        slide_13_mlops, slide_14_limites, slide_15_perspectives, slide_16_conclusion,
        slide_17_merci,
    ]:
        builder(prs)

    prs.save(OUTPUT_PATH)
    print(f"Saved {OUTPUT_PATH} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
